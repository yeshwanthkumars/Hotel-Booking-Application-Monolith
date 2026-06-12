from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette ──────────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT_BLUE  = RGBColor(0x00, 0x8C, 0xFF)   # bright blue
ACCENT_GREEN = RGBColor(0x00, 0xE0, 0x96)   # teal-green
ACCENT_GOLD  = RGBColor(0xFF, 0xC1, 0x07)   # gold
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xD6, 0xE8)
DARK_CARD    = RGBColor(0x1A, 0x2D, 0x45)   # card bg

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

DIAGRAM_BASE = r"C:\Users\YeshwanthKumarS\IntelliJ Spring Boot Workspace - Copilot\Simple-Hotel-Booking-App\docs\diagrams"

# ─────────────────────────────────────────────────────────────────────────────
# Helper utilities
# ─────────────────────────────────────────────────────────────────────────────
def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def set_bg(slide, color=BG_DARK):
    from pptx.oxml.ns import qn
    from lxml import etree
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_multiline(slide, lines, l, t, w, h, font_size=14, color=WHITE,
                  bold_first=False, line_spacing=None):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)
    return txBox

def add_image_safe(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        # placeholder box
        add_rect(slide, l, t, w, h, DARK_CARD)
        add_text_box(slide, f"[Image: {os.path.basename(path)}]",
                     l+0.1, t+h/2-0.2, w-0.2, 0.4,
                     font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

def slide_header(slide, title, subtitle=None):
    # top accent bar
    add_rect(slide, 0, 0, 13.33, 0.08, ACCENT_BLUE)
    add_text_box(slide, title, 0.4, 0.15, 12.5, 0.65,
                 font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle, 0.4, 0.78, 12.5, 0.35,
                     font_size=14, color=ACCENT_GREEN)
    # bottom accent bar
    add_rect(slide, 0, 7.38, 13.33, 0.12, ACCENT_BLUE)

def card(slide, l, t, w, h, title, lines, title_color=ACCENT_GREEN, font_size=13):
    add_rect(slide, l, t, w, h, DARK_CARD)
    add_rect(slide, l, t, w, 0.04, title_color)          # top stripe
    add_text_box(slide, title, l+0.15, t+0.1, w-0.3, 0.4,
                 font_size=14, bold=True, color=title_color)
    add_multiline(slide, lines, l+0.15, t+0.55, w-0.3, h-0.65,
                  font_size=font_size, color=WHITE)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title (Interactive / Visual)
# ═════════════════════════════════════════════════════════════════════════════
s1 = blank_slide(prs)
set_bg(s1)

# ── Left accent strip
add_rect(s1, 0, 0, 0.18, 7.5, ACCENT_BLUE)

# ── Top-right decorative corner block
add_rect(s1, 10.5, 0, 2.83, 0.12, ACCENT_GREEN)
add_rect(s1, 12.5, 0, 0.83, 2.2, DARK_CARD)

# ── Bottom bar
add_rect(s1, 0.18, 7.2, 13.15, 0.3, ACCENT_BLUE)

# ── Big hotel icon area (left panel)
add_rect(s1, 0.18, 0, 4.2, 7.2, DARK_CARD)
add_text_box(s1, "🏨", 0.3, 1.0, 4.0, 1.6, font_size=90, align=PP_ALIGN.CENTER)
add_text_box(s1, "AI-POWERED\nHOTEL BOOKING", 0.3, 2.7, 4.0, 1.3,
             font_size=18, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_rect(s1, 0.7, 4.1, 3.2, 0.05, ACCENT_GREEN)

# Presenter box inside left panel
add_rect(s1, 0.3, 4.3, 3.9, 1.35, RGBColor(0x0A, 0x22, 0x38))
add_text_box(s1, "PRESENTER", 0.4, 4.35, 3.7, 0.35,
             font_size=10, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
add_text_box(s1, "Yeshwanth Kumar S", 0.4, 4.7, 3.7, 0.5,
             font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s1, "June 2026", 0.4, 5.22, 3.7, 0.35,
             font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ── Right content area
# Main title
add_text_box(s1, "Hotel Booking App", 4.65, 0.35, 8.4, 1.1,
             font_size=40, bold=True, color=WHITE)
add_text_box(s1, "Built with AI Assistance", 4.65, 1.42, 8.4, 0.55,
             font_size=22, bold=False, color=ACCENT_GREEN)
add_text_box(s1, "ChatGPT  ·  GitHub Copilot  ·  Anthropic Claude", 4.65, 1.95, 8.4, 0.45,
             font_size=14, color=LIGHT_GRAY)

add_rect(s1, 4.65, 2.52, 8.4, 0.05, ACCENT_GOLD)

# Use case statement
add_rect(s1, 4.65, 2.68, 8.4, 0.95, RGBColor(0x0A, 0x22, 0x38))
add_text_box(s1,
             "💡  Use Case: A complete hotel booking system where users can browse hotels,\n"
             "      book rooms, process payments, cancel bookings & search using natural language AI",
             4.75, 2.72, 8.2, 0.85, font_size=13, color=ACCENT_GOLD, italic=True)

# ── 5 Stats boxes (2 rows)
stats = [
    ("⏱️",  "~7 Days",         "Built with AI",       ACCENT_BLUE),
    ("📡",  "30+ APIs",        "Spring Boot REST",    ACCENT_GREEN),
    ("⚛️",  "Full Stack",      "React + Spring Boot", ACCENT_BLUE),
    ("🤖",  "AI Search",       "Claude + LangChain4j",ACCENT_GOLD),
    ("✅",  "82 Tests",        "Zero Failures",       ACCENT_GREEN),
]
box_w = 1.55
box_h = 1.05
y_stats = 3.82
for i, (icon, val, sub, color) in enumerate(stats):
    x = 4.65 + i * 1.64
    add_rect(s1, x, y_stats, box_w, box_h, DARK_CARD)
    add_rect(s1, x, y_stats, box_w, 0.05, color)
    add_text_box(s1, icon, x, y_stats + 0.04, box_w, 0.38,
                 font_size=20, align=PP_ALIGN.CENTER)
    add_text_box(s1, val, x, y_stats + 0.4, box_w, 0.32,
                 font_size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text_box(s1, sub, x, y_stats + 0.7, box_w, 0.3,
                 font_size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ── Tech stack strip
add_rect(s1, 4.65, 5.05, 8.4, 0.62, RGBColor(0x0A, 0x22, 0x38))
add_text_box(s1, "TECH STACK", 4.75, 5.08, 2.0, 0.28,
             font_size=9, bold=True, color=ACCENT_GREEN)
techs = ["Spring Boot 3", "Java 21", "React", "H2 Database",
         "JWT Security", "LangChain4j", "Swagger UI", "Maven"]
tech_str = "  ·  ".join(techs)
add_text_box(s1, tech_str, 4.75, 5.35, 8.2, 0.28,
             font_size=11, color=WHITE)

# ── Feature highlights row
features = [
    ("🔐", "JWT Auth\n& Roles"),
    ("🏨", "Hotel &\nRoom CRUD"),
    ("💳", "Mock\nPayment"),
    ("❌", "Cancel &\nAuto-Refund"),
    ("🤖", "AI Natural\nSearch"),
]
y_feat = 5.85
fw = 1.55
for i, (icon, label) in enumerate(features):
    x = 4.65 + i * 1.64
    add_rect(s1, x, y_feat, fw, 1.12, RGBColor(0x06, 0x1A, 0x2E))
    add_rect(s1, x, y_feat + 1.07, fw, 0.05, ACCENT_GREEN)
    add_text_box(s1, icon, x, y_feat + 0.05, fw, 0.4,
                 font_size=18, align=PP_ALIGN.CENTER)
    add_text_box(s1, label, x, y_feat + 0.48, fw, 0.6,
                 font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — How I Used AI (The Strategy)
# ═════════════════════════════════════════════════════════════════════════════
s2 = blank_slide(prs)
set_bg(s2)
slide_header(s2, "How I Used AI to Build This", "Two tools. Two different jobs.")

# ChatGPT card
card(s2, 0.4, 1.2, 5.8, 4.5, "🧠  ChatGPT  —  Architect",
     ["Meta Prompting & Design",
      "",
      "✦  Designed full architecture upfront",
      "✦  Generated prompts for GitHub Copilot",
      "✦  Planned features & DB schema",
      "✦  Wrote system design decisions",
      "✦  Drafted user stories & API contracts"],
     title_color=ACCENT_GOLD)

# Copilot card
card(s2, 6.8, 1.2, 5.8, 4.5, "⚡  GitHub Copilot  —  Developer",
     ["Code Generation (IntelliJ + VS Code)",
      "",
      "✦  Generated every backend file",
      "✦  Followed existing coding standards",
      "✦  Wrote 82 unit tests automatically",
      "✦  Generated React UI components",
      "✦  Swagger annotations & configs"],
     title_color=ACCENT_BLUE)

# Arrow between
add_text_box(s2, "→", 6.1, 3.1, 0.6, 0.6, font_size=28, bold=True,
             color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

add_text_box(s2,
             "Workflow:   ChatGPT  →  Design Prompt  →  GitHub Copilot  →  Working Code",
             0.4, 5.9, 12.5, 0.45, font_size=13, bold=True,
             color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What I Built (Full Stack Overview)
# ═════════════════════════════════════════════════════════════════════════════
s3 = blank_slide(prs)
set_bg(s3)
slide_header(s3, "What I Built", "A Full-Stack Monolithic Application")

# Three architecture columns
cols = [
    ("⚛️  React Frontend", ACCENT_BLUE,
     ["• Hotel listing + filters",
      "• Room browsing page",
      "• Booking form & flow",
      "• Payment UI",
      "• Cancel booking",
      "• AI Search box",
      "• JWT token management",
      "• Responsive design"]),
    ("🌱  Spring Boot Backend", ACCENT_GREEN,
     ["• 30+ REST APIs",
      "• JWT Authentication",
      "• Role-based access",
      "• Mock Payment module",
      "• Booking management",
      "• AI Natural Language Search",
      "• Global error handling",
      "• Actuator + Swagger"]),
    ("🗄️  H2 Database", ACCENT_GOLD,
     ["• In-Memory H2 DB",
      "• JPA / Hibernate ORM",
      "• 4 core tables:",
      "  users, hotels,",
      "  rooms, bookings",
      "• Auto-seeded data",
      "• H2 Console UI",
      "• DDL auto create-drop"]),
]
for i, (title, color, lines) in enumerate(cols):
    x = 0.4 + i * 4.3
    card(s3, x, 1.2, 4.0, 5.5, title, lines, title_color=color, font_size=13)

add_text_box(s3, "◀─────── HTTP/REST + JWT ────────▶       ◀──── JPA/SQL ────▶",
             0.4, 5.95, 12.5, 0.4, font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Application Features / Use Cases
# ═════════════════════════════════════════════════════════════════════════════
s4 = blank_slide(prs)
set_bg(s4)
slide_header(s4, "Application Features & Use Cases", "What USER and ADMIN can do")

card(s4, 0.4, 1.2, 5.9, 5.5, "👤  USER Role",
     ["✅  Register & Login  (JWT token)",
      "✅  Browse hotels  (filter: city, stars, price)",
      "✅  Browse available rooms",
      "✅  Create a booking  →  PENDING / UNPAID",
      "✅  Process payment  →  CONFIRMED / PAID",
      "      Gets mock TXN-XXXXXXXX reference",
      "✅  Cancel booking  →  auto-refund if PAID",
      "✅  AI natural language hotel search"],
     title_color=ACCENT_BLUE, font_size=14)

card(s4, 7.0, 1.2, 5.9, 5.5, "🔑  ADMIN Role",
     ["✅  Everything USER can do  +",
      "✅  Add / Edit / Delete hotels",
      "✅  Add / Edit / Delete rooms",
      "✅  View ALL bookings",
      "✅  Refund any booking",
      "✅  Manage all users' data",
      "✅  Full Swagger API access",
      "✅  H2 Console DB viewer"],
     title_color=ACCENT_GOLD, font_size=14)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — E2E Diagram
# ═════════════════════════════════════════════════════════════════════════════
s5 = blank_slide(prs)
set_bg(s5)
slide_header(s5, "End-to-End Application Flow", "Login → Book → Pay → Cancel → AI Search")

e2e_path = os.path.join(DIAGRAM_BASE, "e2e", "HotelBookingCompleteFlow.png")
add_image_safe(s5, e2e_path, 0.4, 1.1, 12.5, 6.0)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — How Efficiently AI Was Used
# ═════════════════════════════════════════════════════════════════════════════
s6 = blank_slide(prs)
set_bg(s6)
slide_header(s6, "How Efficiently AI Was Used", "GitHub Copilot generated the majority of the codebase")

card(s6, 0.4, 1.2, 5.9, 4.2, "⚙️  Backend — Copilot Generated",
     ["• All entity classes & enums",
      "• All DTOs, mappers & records",
      "• All service layer methods",
      "• Repository interfaces (JPA)",
      "• JWT security configuration",
      "• Global exception handler",
      "• 82 unit tests (0 failures)",
      "• Swagger / OpenAPI annotations",
      "• application.properties"],
     title_color=ACCENT_GREEN, font_size=13)

card(s6, 7.0, 1.2, 5.9, 4.2, "🎨  Frontend — Copilot Generated",
     ["• Hotel listing & filter page",
      "• Room browsing components",
      "• Booking form & flow",
      "• Payment UI component",
      "• Cancel booking feature",
      "• AI Search box component",
      "• Axios API service layer",
      "• JWT token handling",
      "• Responsive layout"],
     title_color=ACCENT_BLUE, font_size=13)

add_rect(s6, 0.4, 5.65, 12.5, 0.85, DARK_CARD)
add_text_box(s6,
             "✏️  What I did:   Write the prompts  |  Review & adjust code  |  Connect frontend ↔ backend  |  E2E testing",
             0.6, 5.72, 12.1, 0.7, font_size=14, bold=True,
             color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — AI vs Manual Time Comparison
# ═════════════════════════════════════════════════════════════════════════════
s7 = blank_slide(prs)
set_bg(s7)
slide_header(s7, "AI vs Manual: Time Comparison", "Estimated effort for this project")

rows = [
    ("Feature",                   "Without AI",  "With AI (Copilot)", True),
    ("Architecture & Design",     "2 days",       "4 hours",          False),
    ("Backend APIs  (30+)",       "10 days",      "2 days",           False),
    ("Security  (JWT + Roles)",   "3 days",       "4 hours",          False),
    ("Unit Tests  (82 tests)",    "4 days",       "1 day",            False),
    ("React Frontend",            "7 days",       "2 days",           False),
    ("AI Search Feature",         "5 days",       "1 day",            False),
    ("Docs & Diagrams",           "2 days",       "2 hours",          False),
    ("TOTAL",                     "~33 days",     "~7 days",          True),
]

col_widths = [5.2, 2.8, 3.2]
col_starts = [0.4, 5.8, 8.8]
row_h = 0.48
y_start = 1.15

for r, (feat, wo, wi, is_header) in enumerate(rows):
    y = y_start + r * row_h
    bg = DARK_CARD if not is_header else RGBColor(0x0A, 0x33, 0x55)
    if r == 0:
        bg = ACCENT_BLUE
    add_rect(s7, col_starts[0]-0.05, y, sum(col_widths)+0.3, row_h-0.04, bg)

    fc_feat  = WHITE       if not is_header else ACCENT_GOLD
    fc_wo    = RGBColor(0xFF, 0x77, 0x77)
    fc_wi    = ACCENT_GREEN
    if r == 0:
        fc_feat = fc_wo = fc_wi = WHITE

    fs = 13 if not is_header else 14
    bold = is_header or r == 0
    add_text_box(s7, feat, col_starts[0], y+0.06, col_widths[0], row_h-0.08,
                 font_size=fs, bold=bold, color=fc_feat)
    add_text_box(s7, wo,   col_starts[1], y+0.06, col_widths[1], row_h-0.08,
                 font_size=fs, bold=bold, color=fc_wo, align=PP_ALIGN.CENTER)
    add_text_box(s7, wi,   col_starts[2], y+0.06, col_widths[2], row_h-0.08,
                 font_size=fs, bold=bold, color=fc_wi, align=PP_ALIGN.CENTER)

add_rect(s7, 0.4, 6.3, 12.5, 0.75, DARK_CARD)
add_text_box(s7,
             '🚀  Productivity gain:  ~5x faster with AI assistance\n'
             '"AI didn\'t write perfect code — it wrote the first draft fast. I reviewed, refined, and connected the pieces."',
             0.6, 6.33, 12.1, 0.68, font_size=13, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Sample Prompts Used
# ═════════════════════════════════════════════════════════════════════════════
s8 = blank_slide(prs)
set_bg(s8)
slide_header(s8, "Sample Prompts Used", "Meta prompting with ChatGPT → implementation with GitHub Copilot")

prompts = [
    ("🧠  Meta Prompt  (ChatGPT — Architecture)", ACCENT_GOLD,
     '"Design a Spring Boot REST API for hotel booking with JWT auth,\n'
     ' role-based access (USER/ADMIN), entities for Hotel/Room/Booking/User,\n'
     ' booking overlap validation, and pagination. Give me the full\n'
     ' architecture and a GitHub Copilot prompt to implement it."'),
    ("⚡  Feature Prompt  (GitHub Copilot — Cancellation)", ACCENT_BLUE,
     '"Implement BookingService.cancelBooking(id, username, isAdmin).\n'
     ' USER can cancel own booking. ADMIN can cancel any.\n'
     ' If booking was PAID, auto-set paymentStatus=REFUNDED.\n'
     ' Return CancellationResponse record. Follow existing patterns. No new files."'),
    ("🤖  AI Search Prompt  (GitHub Copilot — LangChain4j)", ACCENT_GREEN,
     '"Add AI natural language hotel search using LangChain4j + Claude.\n'
     ' Parse query → extract city, starRating, dates, price range.\n'
     ' Phase 1: search hotels. Phase 2: check room availability.\n'
     ' Reuse HotelService, RoomService, BookingRepository. Zero new DB queries."'),
]

for i, (title, color, prompt_text) in enumerate(prompts):
    y = 1.2 + i * 1.95
    add_rect(s8, 0.4, y, 12.5, 1.82, DARK_CARD)
    add_rect(s8, 0.4, y, 12.5, 0.04, color)
    add_text_box(s8, title, 0.55, y+0.08, 5.0, 0.38,
                 font_size=13, bold=True, color=color)
    add_text_box(s8, prompt_text, 0.55, y+0.45, 12.1, 1.28,
                 font_size=12, color=LIGHT_GRAY, italic=True)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — AI Search Feature Deep Dive
# ═════════════════════════════════════════════════════════════════════════════
s9 = blank_slide(prs)
set_bg(s9)
slide_header(s9, "AI Search — The Highlight Feature ✨",
             "Natural Language Hotel Search using LangChain4j + Anthropic Claude")

# Flow boxes
flow_steps = [
    ('User Types\nNatural Language Query', ACCENT_BLUE),
    ('Claude AI\nExtracts Parameters', ACCENT_GOLD),
    ('Phase 1\nHotel Search', ACCENT_GREEN),
    ('Phase 2\nAvailability Check', ACCENT_GREEN),
    ('Results\nReturned', ACCENT_BLUE),
]
box_w = 2.1
box_h = 0.85
y_box = 1.3
for i, (label, color) in enumerate(flow_steps):
    x = 0.35 + i * 2.6
    add_rect(s9, x, y_box, box_w, box_h, color)
    add_text_box(s9, label, x, y_box+0.05, box_w, box_h-0.1,
                 font_size=12, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    if i < len(flow_steps) - 1:
        add_text_box(s9, "→", x+box_w+0.02, y_box+0.2, 0.45, 0.45,
                     font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Example query
add_rect(s9, 0.4, 2.35, 12.5, 0.55, DARK_CARD)
add_text_box(s9,
             '💬  Query:  "5 star hotels in Dubai available August 10 to August 15"',
             0.6, 2.4, 12.1, 0.45, font_size=14, color=ACCENT_GOLD, italic=True)

# Claude extracts
add_rect(s9, 0.4, 3.05, 5.9, 1.85, DARK_CARD)
add_rect(s9, 0.4, 3.05, 5.9, 0.04, ACCENT_GOLD)
add_text_box(s9, "🧠  Claude Extracts (JSON)", 0.55, 3.1, 5.6, 0.4,
             font_size=13, bold=True, color=ACCENT_GOLD)
add_text_box(s9,
             '{ "city": "Dubai",\n  "starRating": 5,\n  "checkIn": "2026-08-10",\n  "checkOut": "2026-08-15" }',
             0.55, 3.55, 5.6, 1.3, font_size=12, color=ACCENT_GREEN, italic=True)

# Why efficient
card(s9, 6.8, 3.05, 5.9, 1.85, "⚡  Why It's Efficient",
     ["✅  150 lines of new code total",
      "✅  0 existing lines changed",
      "✅  0 new DB tables created",
      "✅  Reuses HotelService, RoomService,",
      "     BookingRepository — unchanged",
      "✅  Understands relative dates"],
     title_color=ACCENT_GREEN, font_size=13)

# Graceful degradation
add_rect(s9, 0.4, 5.1, 12.5, 1.3, RGBColor(0x0A, 0x22, 0x38))
add_text_box(s9, "🛡️  Graceful Degradation", 0.6, 5.15, 5.0, 0.38,
             font_size=13, bold=True, color=ACCENT_BLUE)
add_text_box(s9,
             "No dates in query → Phase 2 skipped, hotels still returned   |   "
             "Claude API down → 500 with meaningful message   |   "
             "Vague query → returns all hotels",
             0.6, 5.55, 12.1, 0.75, font_size=12, color=LIGHT_GRAY)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Demo + Key Takeaways
# ═════════════════════════════════════════════════════════════════════════════
s10 = blank_slide(prs)
set_bg(s10)
slide_header(s10, "Live Demo & Key Takeaways", "")

card(s10, 0.4, 1.1, 6.0, 5.5, "🎬  Demo Flow",
     ["1️⃣   Login as USER / ADMIN  (get JWT)",
      "2️⃣   Browse Hotels  (filter: city, stars)",
      "3️⃣   Browse Rooms  (filter: hotel, type)",
      "4️⃣   Create Booking  →  PENDING / UNPAID",
      "5️⃣   Process Payment  →  CONFIRMED / PAID",
      "        TXN-XXXXXXXX reference generated",
      "6️⃣   Cancel Booking  →  CANCELLED / REFUNDED",
      "7️⃣   AI Search: \"luxury hotels in Singapore\"",
      "8️⃣   Swagger UI  →  all 30+ APIs live",
      "9️⃣   H2 Console  →  live booking table"],
     title_color=ACCENT_BLUE, font_size=13)

card(s10, 7.0, 1.1, 5.9, 2.55, "💡  Key Takeaways",
     ["🧠  ChatGPT   =  Architect  (design & plan)",
      "⚡  Copilot   =  Developer  (code & tests)",
      "🤖  Claude    =  Feature    (AI intelligence)",
      "👤  You       =  Tech Lead  (review & decide)"],
     title_color=ACCENT_GOLD, font_size=14)

card(s10, 7.0, 3.85, 5.9, 2.75, "🔗  URLs",
     ["App:      http://localhost:8080",
      "Swagger:  http://localhost:8080/swagger-ui.html",
      "H2 DB:    http://localhost:8080/h2-console",
      "AI API:   POST /api/v1/ai/search",
      "Creds:    user/password  |  admin/password"],
     title_color=ACCENT_GREEN, font_size=13)

add_rect(s10, 0.4, 6.8, 12.5, 0.55, DARK_CARD)
add_text_box(s10,
             '"AI tools don\'t replace developers — they make developers 5x more productive."',
             0.6, 6.84, 12.1, 0.46, font_size=14, bold=True,
             color=ACCENT_GOLD, align=PP_ALIGN.CENTER, italic=True)

# ═════════════════════════════════════════════════════════════════════════════
# Save
# ═════════════════════════════════════════════════════════════════════════════
out_path = r"C:\Users\YeshwanthKumarS\IntelliJ Spring Boot Workspace - Copilot\Simple-Hotel-Booking-App\Hotel_Booking_App_AI_POC.pptx"
prs.save(out_path)
print(f"✅  PPT saved: {out_path}")

